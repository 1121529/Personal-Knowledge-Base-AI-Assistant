"""
parser.py

RAG 專用文件解析器

支援：
- PDF
- DOCX
- PPTX

輸出格式：

[
    {
        "source": "numpy.pdf",
        "page": 1,
        "content": "..."
    }
]

Author: Team RAG
"""

from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation


class DocumentParser:

    @staticmethod
    def parse_pdf(file_path: str):
        """
        PDF解析
        """

        documents = []

        source = Path(file_path).name

        with pdfplumber.open(file_path) as pdf:

            for page_num, page in enumerate(pdf.pages, start=1):

                text = page.extract_text()

                if text and text.strip():

                    documents.append({
                        "source": source,
                        "page": page_num,
                        "content": text.strip()
                    })

        return documents

    @staticmethod
    def parse_docx(file_path: str):
        """
        DOCX解析
        """

        source = Path(file_path).name

        doc = Document(file_path)

        text_list = []

        for paragraph in doc.paragraphs:

            text = paragraph.text.strip()

            if text:
                text_list.append(text)

        return [{
            "source": source,
            "page": 1,
            "content": "\n".join(text_list)
        }]

    @staticmethod
    def parse_pptx(file_path: str):
        """
        PPTX解析
        """

        source = Path(file_path).name

        prs = Presentation(file_path)

        documents = []

        for slide_num, slide in enumerate(prs.slides, start=1):

            slide_text = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:
                        slide_text.append(text)

            if slide_text:

                documents.append({
                    "source": source,
                    "page": slide_num,
                    "content": "\n".join(slide_text)
                })

        return documents

    @staticmethod
    def parse(file_path: str):
        """
        自動判斷格式
        """

        suffix = Path(file_path).suffix.lower()

        if suffix == ".pdf":
            return DocumentParser.parse_pdf(file_path)

        elif suffix == ".docx":
            return DocumentParser.parse_docx(file_path)

        elif suffix == ".pptx":
            return DocumentParser.parse_pptx(file_path)

        else:
            raise ValueError(
                f"不支援的檔案格式: {suffix}"
            )


if __name__ == "__main__":

    file_path = "uploads/numpy.pdf"

    docs = DocumentParser.parse(file_path)

    print(f"解析頁數: {len(docs)}")

    print("\n第一筆資料:\n")

    print(docs[0])
